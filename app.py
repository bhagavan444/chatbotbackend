from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import google.generativeai as genai

import os
import uuid
from datetime import datetime
from threading import Lock

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import logging

# ---------------- ENV ----------------
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    raise RuntimeError("❌ GEMINI_API_KEY not found in environment variables")

# ---------------- LOGGING ----------------
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ---------------- GEMINI CONFIG ----------------
genai.configure(api_key=GEMINI_API_KEY)

# ---------------- LAZY MODEL LOADER (DEPLOY SAFE) ----------------
_model = None
_model_lock = Lock()

def get_working_model():
    global _model
    with _model_lock:
        if _model is not None:
            return _model

        for m in genai.list_models():
            if "generateContent" in m.supported_generation_methods:
                logger.info(f"✅ Using Gemini model: {m.name}")
                _model = genai.GenerativeModel(m.name)
                return _model

        raise RuntimeError("❌ No supported Gemini model found")

# ---------------- APP SETUP ----------------
app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True)

# ---------------- STORAGE ----------------
chat_sessions = {}
chat_lock = Lock()

# ---------------- FILE STORAGE ----------------
DOWNLOAD_DIR = "downloads"
os.makedirs(DOWNLOAD_DIR, exist_ok=True)

# ---------------- FILE TEXT EXTRACTION ----------------
def extract_file_text(path):
    ext = os.path.splitext(path)[1].lower()
    text = ""

    try:
        if ext == ".pdf":
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)

        elif ext == ".docx":
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs)

        elif ext == ".pptx":
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        logger.warning(f"File parse error: {e}")

    return text.strip()

# ---------------- GEMINI CALL ----------------
def call_gemini(prompt):
    try:
        model = get_working_model()
        response = model.generate_content(prompt)

        if response and response.candidates:
            candidate = response.candidates[0]
            if candidate.content.parts:
                return candidate.content.parts[0].text

        return "⚠️ No response from AI."

    except Exception as e:
        logger.error(f"Gemini error: {e}")
        return "⚠️ Gemini service temporarily unavailable."

# ---------------- CHAT API ----------------
@app.route("/api/chat", methods=["POST"])
def chat():
    try:
        user_input = ""
        chat_id = None
        files = []

        if request.content_type and request.content_type.startswith("multipart/form-data"):
            user_input = request.form.get("message", "")
            chat_id = request.form.get("chat_id")
            files = request.files.getlist("files")
        else:
            data = request.get_json() or {}
            user_input = data.get("message", "")
            chat_id = data.get("chat_id")

        if not user_input and not files:
            return jsonify({"reply": "⚠️ No input provided"}), 400

        uploaded_files = []
        file_text = ""

        for file in files:
            filename = f"{uuid.uuid4().hex}_{file.filename}"
            path = os.path.join(DOWNLOAD_DIR, filename)
            file.save(path)
            uploaded_files.append(filename)
            file_text += "\n" + extract_file_text(path)

        final_input = user_input + "\n\n" + file_text if file_text else user_input

        prompt = f"""
You are a helpful AI assistant.
Answer clearly in your own words.

User input:
{final_input}
"""

        reply = call_gemini(prompt)

        new_chat_id = chat_id or str(uuid.uuid4())
        timestamp = datetime.now().strftime("%H:%M")

        with chat_lock:
            chat_sessions.setdefault(new_chat_id, {"messages": []})

            chat_sessions[new_chat_id]["messages"].append({
                "id": str(uuid.uuid4()),
                "role": "user",
                "message": user_input,
                "time": timestamp,
                "files": uploaded_files or None
            })

            chat_sessions[new_chat_id]["messages"].append({
                "id": str(uuid.uuid4()),
                "role": "assistant",
                "reply": reply,
                "time": timestamp
            })

        return jsonify({"reply": reply, "chat_id": new_chat_id})

    except Exception:
        logger.exception("Chat error")
        return jsonify({"reply": "Server error"}), 500

# ---------------- CHAT HISTORY ----------------
@app.route("/api/chats", methods=["GET"])
def get_chats():
    return jsonify(chat_sessions)

# ---------------- FILE DOWNLOAD ----------------
@app.route("/download/<filename>")
def download(filename):
    return send_from_directory(DOWNLOAD_DIR, filename, as_attachment=True)

# ---------------- RUN ----------------
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
